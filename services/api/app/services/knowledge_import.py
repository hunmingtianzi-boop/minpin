from __future__ import annotations

import csv
import html
import io
import json
import re
import zipfile
from dataclasses import dataclass
from functools import lru_cache
from html.parser import HTMLParser
from pathlib import PurePath
from typing import Literal

from docx import Document
from openpyxl import load_workbook
from PIL import Image, UnidentifiedImageError
from pptx import Presentation
from pypdf import PdfReader

MAX_FILE_BYTES = 10 * 1024 * 1024
MAX_BATCH_BYTES = 25 * 1024 * 1024
MAX_FILES = 5
MAX_CSV_ROWS = 1_000
MAX_CSV_COLUMNS = 16
MAX_CSV_CELL_CHARS = 100_000
MAX_TEXT_CHARS = 1_000_000
MAX_PDF_PAGES = 500
MAX_OCR_PAGES = 50
MAX_IMAGE_PIXELS = 20_000_000
MAX_PPTX_SLIDES = 500
MAX_XLSX_SHEETS = 100
MAX_ARCHIVE_ENTRIES = 2_000
MAX_ARCHIVE_UNCOMPRESSED_BYTES = 50 * 1024 * 1024
MAX_ARCHIVE_ENTRY_BYTES = 20 * 1024 * 1024
_DANGEROUS_PREFIXES = ("=", "+", "-", "@", "\t", "\r")
_CONTROL_RE = re.compile(r"[\x00-\x08\x0b\x0c\x0e-\x1f]")
_MIMES = {
    "pdf": {"application/pdf"},
    "docx": {"application/vnd.openxmlformats-officedocument.wordprocessingml.document"},
    "csv": {"text/csv", "application/csv", "text/plain", "application/vnd.ms-excel"},
    "pptx": {"application/vnd.openxmlformats-officedocument.presentationml.presentation"},
    "xlsx": {"application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"},
    "txt": {"text/plain"},
    "md": {"text/markdown", "text/plain"},
    "html": {"text/html", "application/xhtml+xml"},
    "htm": {"text/html", "application/xhtml+xml"},
    "png": {"image/png"},
    "jpg": {"image/jpeg"},
    "jpeg": {"image/jpeg"},
    "webp": {"image/webp"},
    "tiff": {"image/tiff"},
    "bmp": {"image/bmp"},
}
_IMAGE_TYPES = frozenset({"png", "jpg", "jpeg", "webp", "tiff", "bmp"})


class KnowledgeImportError(ValueError):
    def __init__(self, code: str) -> None:
        super().__init__(code)
        self.code = code


@dataclass(frozen=True, slots=True)
class ImportDraft:
    title: str
    raw_text: str
    visibility: Literal["public", "authenticated", "internal"]
    auto_publish: bool = False


def safe_file_name(value: str | None) -> str:
    name = (value or "").strip()
    if not name or name in {".", ".."} or PurePath(name).name != name:
        raise KnowledgeImportError("IMPORT_UNSAFE_FILENAME")
    if any(character in name for character in ("/", "\\", "\x00")):
        raise KnowledgeImportError("IMPORT_UNSAFE_FILENAME")
    return name[:255]


def validate_upload(name: str, content_type: str | None, payload: bytes) -> str:
    if not payload:
        raise KnowledgeImportError("IMPORT_EMPTY_FILE")
    if len(payload) > MAX_FILE_BYTES:
        raise KnowledgeImportError("IMPORT_FILE_TOO_LARGE")
    extension = name.rsplit(".", 1)[-1].casefold() if "." in name else ""
    if extension not in _MIMES:
        raise KnowledgeImportError("IMPORT_UNSUPPORTED_TYPE")
    normalized_mime = (content_type or "").split(";", 1)[0].strip().casefold()
    if normalized_mime not in _MIMES[extension]:
        raise KnowledgeImportError("IMPORT_MIME_MISMATCH")
    if extension == "pdf" and not payload.startswith(b"%PDF-"):
        raise KnowledgeImportError("IMPORT_MAGIC_MISMATCH")
    if extension == "docx":
        if not payload.startswith(b"PK\x03\x04"):
            raise KnowledgeImportError("IMPORT_MAGIC_MISMATCH")
        _validate_office_archive(
            payload, required_entry="word/document.xml", error="IMPORT_DOCX_INVALID"
        )
    if extension == "pptx":
        if not payload.startswith(b"PK\x03\x04"):
            raise KnowledgeImportError("IMPORT_MAGIC_MISMATCH")
        _validate_office_archive(
            payload, required_entry="ppt/presentation.xml", error="IMPORT_PPTX_INVALID"
        )
    if extension == "xlsx":
        if not payload.startswith(b"PK\x03\x04"):
            raise KnowledgeImportError("IMPORT_MAGIC_MISMATCH")
        _validate_office_archive(
            payload, required_entry="xl/workbook.xml", error="IMPORT_XLSX_INVALID"
        )
    if extension == "csv" and (payload.startswith(b"%PDF-") or payload.startswith(b"PK\x03\x04")):
        raise KnowledgeImportError("IMPORT_MAGIC_MISMATCH")
    if extension in _IMAGE_TYPES:
        _validate_image(payload)
    return extension


def parse_payload(source_type: str, file_name: str, payload: bytes) -> list[ImportDraft]:
    if source_type == "pdf":
        return [_parse_pdf(file_name, payload)]
    if source_type == "docx":
        return [_parse_docx(file_name, payload)]
    if source_type == "csv":
        return [_parse_csv(file_name, payload)]
    if source_type == "pptx":
        return [_parse_pptx(file_name, payload)]
    if source_type == "xlsx":
        return [_parse_xlsx(file_name, payload)]
    if source_type in {"txt", "md", "html", "htm"}:
        return [_parse_text(file_name, source_type, payload)]
    if source_type in _IMAGE_TYPES:
        return [_parse_image(file_name, payload)]
    raise KnowledgeImportError("IMPORT_UNSUPPORTED_TYPE")


def encode_draft(draft: ImportDraft) -> bytes:
    return json.dumps(
        {
            "title": draft.title,
            "raw_text": draft.raw_text,
            "visibility": draft.visibility,
            "auto_publish": draft.auto_publish,
        },
        ensure_ascii=False,
        separators=(",", ":"),
    ).encode("utf-8")


def decode_draft(payload: bytes) -> ImportDraft:
    try:
        value = json.loads(payload.decode("utf-8"))
        return _validated_draft(
            value["title"],
            value["raw_text"],
            value["visibility"],
            auto_publish=bool(value.get("auto_publish", False)),
        )
    except (KeyError, TypeError, ValueError, UnicodeDecodeError, json.JSONDecodeError) as exc:
        raise KnowledgeImportError("IMPORT_PAYLOAD_INVALID") from exc


def _parse_pdf(file_name: str, payload: bytes) -> ImportDraft:
    text = ""
    page_count: int | None = None
    try:
        # Some perfectly viewable PDFs contain non-standard cross-reference
        # tables.  Tolerant parsing keeps those documents importable, while
        # the independent PyMuPDF fallback below still rejects invalid files.
        reader = PdfReader(io.BytesIO(payload), strict=False)
        if reader.is_encrypted:
            raise KnowledgeImportError("IMPORT_ENCRYPTED_PDF")
        page_count = len(reader.pages)
        if page_count > MAX_PDF_PAGES:
            raise KnowledgeImportError("IMPORT_PDF_TOO_MANY_PAGES")
        text = "\n\n".join((page.extract_text() or "").strip() for page in reader.pages).strip()
    except KnowledgeImportError:
        raise
    except Exception:
        # PyMuPDF accepts a number of PDFs that pypdf cannot recover.  Do not
        # fail until both local parsers have rejected the upload.
        page_count = None

    if page_count is None or len(text) < 80:
        try:
            fallback_text, fallback_page_count = _extract_pdf_text_with_fitz(payload)
            page_count = fallback_page_count
            if len(fallback_text) > len(text):
                text = fallback_text
        except KnowledgeImportError:
            if page_count is None:
                raise

    if page_count is None:
        raise KnowledgeImportError("IMPORT_PDF_INVALID")
    if len(text) < 80:
        text = _ocr_pdf(payload, max_pages=min(MAX_OCR_PAGES, page_count))
    return _validated_draft(file_name.rsplit(".", 1)[0], text, "public")


def _extract_pdf_text_with_fitz(payload: bytes) -> tuple[str, int]:
    """Extract selectable text with a parser that is tolerant of damaged xref data."""

    try:
        import fitz

        document = fitz.open(stream=payload, filetype="pdf")
        try:
            if document.needs_pass:
                raise KnowledgeImportError("IMPORT_ENCRYPTED_PDF")
            if document.page_count > MAX_PDF_PAGES:
                raise KnowledgeImportError("IMPORT_PDF_TOO_MANY_PAGES")
            values = [
                (document.load_page(page_number).get_text("text") or "").strip()
                for page_number in range(document.page_count)
            ]
            return "\n\n".join(value for value in values if value).strip(), document.page_count
        finally:
            document.close()
    except KnowledgeImportError:
        raise
    except Exception as exc:
        raise KnowledgeImportError("IMPORT_PDF_INVALID") from exc


def _parse_docx(file_name: str, payload: bytes) -> ImportDraft:
    _validate_office_archive(
        payload, required_entry="word/document.xml", error="IMPORT_DOCX_INVALID"
    )
    try:
        document = Document(io.BytesIO(payload))
        parts = [
            paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()
        ]
        for table in document.tables:
            for row in table.rows:
                line = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if line:
                    parts.append(line)
    except Exception as exc:
        raise KnowledgeImportError("IMPORT_DOCX_INVALID") from exc
    return _validated_draft(file_name.rsplit(".", 1)[0], "\n\n".join(parts), "public")


def _parse_pptx(file_name: str, payload: bytes) -> ImportDraft:
    _validate_office_archive(
        payload, required_entry="ppt/presentation.xml", error="IMPORT_PPTX_INVALID"
    )
    try:
        presentation = Presentation(io.BytesIO(payload))
        if len(presentation.slides) > MAX_PPTX_SLIDES:
            raise KnowledgeImportError("IMPORT_PPTX_TOO_MANY_SLIDES")
        parts: list[str] = []
        for ordinal, slide in enumerate(presentation.slides, start=1):
            values: list[str] = []
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False) and shape.text.strip():
                    values.append(shape.text.strip())
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        line = " | ".join(
                            cell.text.strip() for cell in row.cells if cell.text.strip()
                        )
                        if line:
                            values.append(line)
            if values:
                parts.append(f"第 {ordinal} 页\n" + "\n".join(values))
    except Exception as exc:
        raise KnowledgeImportError("IMPORT_PPTX_INVALID") from exc
    return _validated_draft(file_name.rsplit(".", 1)[0], "\n\n".join(parts), "public")


def _parse_xlsx(file_name: str, payload: bytes) -> ImportDraft:
    _validate_office_archive(payload, required_entry="xl/workbook.xml", error="IMPORT_XLSX_INVALID")
    try:
        workbook = load_workbook(io.BytesIO(payload), read_only=True, data_only=False)
        if len(workbook.worksheets) > MAX_XLSX_SHEETS:
            raise KnowledgeImportError("IMPORT_XLSX_TOO_MANY_SHEETS")
        parts: list[str] = []
        for sheet in workbook.worksheets:
            rows: list[str] = []
            for row in sheet.iter_rows(values_only=True):
                values = [
                    str(value).strip() for value in row if value is not None and str(value).strip()
                ]
                if values:
                    rows.append(" | ".join(values))
                if len(rows) >= MAX_CSV_ROWS:
                    break
            if rows:
                parts.append(f"工作表：{sheet.title}\n" + "\n".join(rows))
    except Exception as exc:
        raise KnowledgeImportError("IMPORT_XLSX_INVALID") from exc
    return _validated_draft(file_name.rsplit(".", 1)[0], "\n\n".join(parts), "public")


def _parse_text(file_name: str, source_type: str, payload: bytes) -> ImportDraft:
    try:
        value = payload.decode("utf-8-sig")
    except UnicodeDecodeError as exc:
        raise KnowledgeImportError("IMPORT_TEXT_ENCODING") from exc
    if source_type in {"html", "htm"}:
        value = _html_to_text(value)
    return _validated_draft(file_name.rsplit(".", 1)[0], value, "public")


def _parse_image(file_name: str, payload: bytes) -> ImportDraft:
    return _validated_draft(file_name.rsplit(".", 1)[0], _ocr_image(payload), "public")


def _parse_csv(file_name: str, payload: bytes) -> ImportDraft:
    try:
        text = payload.decode("utf-8-sig")
    except UnicodeDecodeError as exc:
        raise KnowledgeImportError("IMPORT_CSV_ENCODING") from exc
    if "\x00" in text:
        raise KnowledgeImportError("IMPORT_CSV_INVALID")
    rows: list[str] = []
    visibilities: set[str] = set()
    try:
        reader = csv.DictReader(io.StringIO(text, newline=""))
        raw_fields = [
            str(field).strip() for field in (reader.fieldnames or []) if field is not None
        ]
        fields = set(raw_fields)
        if len(raw_fields) > MAX_CSV_COLUMNS or len(raw_fields) != len(fields):
            raise KnowledgeImportError("IMPORT_CSV_HEADERS")
        if not raw_fields:
            raise KnowledgeImportError("IMPORT_CSV_HEADERS")
        for row in reader:
            if len(rows) >= MAX_CSV_ROWS:
                raise KnowledgeImportError("IMPORT_CSV_TOO_MANY_ROWS")
            if None in row:
                raise KnowledgeImportError("IMPORT_CSV_COLUMNS")
            values = {str(key).strip(): value or "" for key, value in row.items()}
            for value in values.values():
                if len(value) > MAX_CSV_CELL_CHARS:
                    raise KnowledgeImportError("IMPORT_CSV_CELL_TOO_LARGE")
                if value.lstrip().startswith(_DANGEROUS_PREFIXES):
                    raise KnowledgeImportError("IMPORT_DANGEROUS_VALUE")
            visibility = values.get("visibility", "").strip().casefold()
            if visibility:
                if visibility not in {"public", "authenticated", "internal"}:
                    raise KnowledgeImportError("IMPORT_VISIBILITY_INVALID")
                visibilities.add(visibility)
            rows.append(
                " | ".join(
                    f"{field}: {values.get(field, '').strip()}"
                    for field in raw_fields
                    if values.get(field, "").strip()
                )
            )
    except csv.Error as exc:
        raise KnowledgeImportError("IMPORT_CSV_INVALID") from exc
    if not rows:
        raise KnowledgeImportError("IMPORT_EMPTY_TEXT")
    visibility = next(iter(visibilities)) if len(visibilities) == 1 else "public"
    return _validated_draft(
        file_name.rsplit(".", 1)[0], "\n".join([" | ".join(raw_fields), *rows]), visibility
    )


def _validated_draft(
    title: str, raw_text: str, visibility: str, *, auto_publish: bool = False
) -> ImportDraft:
    normalized_title = str(title).strip()
    normalized_text = str(raw_text).strip()
    normalized_visibility = str(visibility).strip().casefold()
    if not normalized_title or len(normalized_title) > 500:
        raise KnowledgeImportError("IMPORT_TITLE_INVALID")
    if not normalized_text:
        raise KnowledgeImportError("IMPORT_EMPTY_TEXT")
    if len(normalized_text) > MAX_TEXT_CHARS:
        raise KnowledgeImportError("IMPORT_TEXT_TOO_LARGE")
    if _CONTROL_RE.search(normalized_title) or _CONTROL_RE.search(normalized_text):
        raise KnowledgeImportError("IMPORT_DANGEROUS_VALUE")
    if normalized_title.lstrip().startswith(_DANGEROUS_PREFIXES):
        raise KnowledgeImportError("IMPORT_DANGEROUS_VALUE")
    if normalized_visibility not in {"public", "authenticated", "internal"}:
        raise KnowledgeImportError("IMPORT_VISIBILITY_INVALID")
    return ImportDraft(  # type: ignore[arg-type]
        normalized_title, normalized_text, normalized_visibility, auto_publish=auto_publish
    )


def _validate_office_archive(payload: bytes, *, required_entry: str, error: str) -> None:
    """Reject archive bombs, traversal entries and macro-enabled Office files.

    OOXML is a ZIP archive.  Validate its central directory before handing it
    to third-party parsers, which prevents a small upload from expanding into
    unbounded disk/memory work in a worker process.
    """
    try:
        with zipfile.ZipFile(io.BytesIO(payload)) as archive:
            infos = archive.infolist()
            if len(infos) > MAX_ARCHIVE_ENTRIES:
                raise KnowledgeImportError("IMPORT_ARCHIVE_LIMIT")
            if sum(item.file_size for item in infos) > MAX_ARCHIVE_UNCOMPRESSED_BYTES:
                raise KnowledgeImportError("IMPORT_ARCHIVE_LIMIT")
            names = [item.filename for item in infos]
            if required_entry not in names:
                raise KnowledgeImportError(error)
            for item in infos:
                name = item.filename
                if item.file_size > MAX_ARCHIVE_ENTRY_BYTES:
                    raise KnowledgeImportError("IMPORT_ARCHIVE_LIMIT")
                if item.file_size > 1_000_000 and item.compress_size * 100 < item.file_size:
                    raise KnowledgeImportError("IMPORT_ARCHIVE_LIMIT")
                path = PurePath(name.replace("\\", "/"))
                if path.is_absolute() or ".." in path.parts:
                    raise KnowledgeImportError("IMPORT_ARCHIVE_PATH")
                lowered = name.casefold()
                if lowered.endswith(("vbaproject.bin", ".docm")) or "macros" in lowered:
                    raise KnowledgeImportError("IMPORT_MACRO_DOCUMENT")
    except KnowledgeImportError:
        raise
    except (OSError, zipfile.BadZipFile) as exc:
        raise KnowledgeImportError(error) from exc


def _validate_image(payload: bytes) -> None:
    try:
        with Image.open(io.BytesIO(payload)) as image:
            image.verify()
        with Image.open(io.BytesIO(payload)) as image:
            width, height = image.size
            if width < 1 or height < 1 or width * height > MAX_IMAGE_PIXELS:
                raise KnowledgeImportError("IMPORT_IMAGE_DIMENSIONS")
    except KnowledgeImportError:
        raise
    except (OSError, UnidentifiedImageError, Image.DecompressionBombError) as exc:
        raise KnowledgeImportError("IMPORT_IMAGE_INVALID") from exc


@lru_cache(maxsize=1)
def _ocr_engine():
    """Load the local OCR engine lazily, only in the background worker."""
    try:
        from rapidocr import RapidOCR
    except ImportError as exc:  # pragma: no cover - exercised in deploy configuration
        raise KnowledgeImportError("IMPORT_OCR_UNAVAILABLE") from exc
    return RapidOCR()


def _ocr_image(payload: bytes) -> str:
    _validate_image(payload)
    try:
        import numpy as np

        with Image.open(io.BytesIO(payload)) as image:
            rgb = image.convert("RGB")
            result = _ocr_engine()(np.asarray(rgb))
        # v3 exposes a result object; keep the tuple/list fallback for a
        # locally pinned v2 engine during rolling upgrades.
        if hasattr(result, "txts"):
            text = "\n".join(str(value).strip() for value in result.txts if str(value).strip())
        else:
            lines = result[0] if isinstance(result, tuple) else result
            text = "\n".join(
                str(line[1]).strip()
                for line in (lines or [])
                if isinstance(line, (list, tuple)) and len(line) > 1 and str(line[1]).strip()
            )
    except KnowledgeImportError:
        raise
    except Exception as exc:
        raise KnowledgeImportError("IMPORT_OCR_FAILED") from exc
    if not text:
        raise KnowledgeImportError("IMPORT_OCR_EMPTY")
    return text


def _ocr_pdf(payload: bytes, *, max_pages: int) -> str:
    try:
        import fitz

        document = fitz.open(stream=payload, filetype="pdf")
        try:
            values: list[str] = []
            errors: list[KnowledgeImportError] = []
            for page_number in range(min(document.page_count, max_pages)):
                try:
                    page = document.load_page(page_number)
                    pixmap = page.get_pixmap(matrix=fitz.Matrix(1.5, 1.5), alpha=False)
                    values.append(_ocr_image(pixmap.tobytes("png")))
                except KnowledgeImportError as exc:
                    # A blank or damaged page must not discard text recognized
                    # from the remaining pages in the same PDF.
                    errors.append(exc)
            text = "\n\n".join(value for value in values if value).strip()
            if text:
                return text
            if errors:
                raise errors[0]
            raise KnowledgeImportError("IMPORT_OCR_EMPTY")
        finally:
            document.close()
    except KnowledgeImportError:
        raise
    except Exception as exc:
        raise KnowledgeImportError("IMPORT_PDF_OCR_FAILED") from exc


class _TextExtractor(HTMLParser):
    _IGNORED = {"script", "style", "noscript", "template"}

    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self._parts: list[str] = []
        self._ignored_depth = 0

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        if tag.casefold() in self._IGNORED:
            self._ignored_depth += 1
        elif tag.casefold() in {"p", "div", "li", "br", "tr", "h1", "h2", "h3", "h4"}:
            self._parts.append("\n")

    def handle_endtag(self, tag: str) -> None:
        if tag.casefold() in self._IGNORED and self._ignored_depth:
            self._ignored_depth -= 1
        elif tag.casefold() in {"p", "div", "li", "tr", "h1", "h2", "h3", "h4"}:
            self._parts.append("\n")

    def handle_data(self, data: str) -> None:
        if not self._ignored_depth:
            self._parts.append(data)

    def text(self) -> str:
        return re.sub(r"\n{3,}", "\n\n", "".join(self._parts)).strip()


def _html_to_text(value: str) -> str:
    parser = _TextExtractor()
    try:
        parser.feed(value)
        parser.close()
    except Exception as exc:
        raise KnowledgeImportError("IMPORT_HTML_INVALID") from exc
    return html.unescape(parser.text())


__all__ = [
    "ImportDraft",
    "KnowledgeImportError",
    "MAX_BATCH_BYTES",
    "MAX_FILES",
    "decode_draft",
    "encode_draft",
    "parse_payload",
    "safe_file_name",
    "validate_upload",
]
