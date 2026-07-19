from __future__ import annotations

import io
import sys
import zipfile
from types import SimpleNamespace

import pytest
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pypdf import PdfWriter

from app.services import knowledge_import as knowledge_import_module
from app.services.knowledge_import import (
    KnowledgeImportError,
    parse_payload,
    safe_file_name,
    validate_upload,
)


def _docx_bytes(text: str) -> bytes:
    document = Document()
    document.add_paragraph(text)
    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def _pptx_bytes(text: str) -> bytes:
    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[1]).shapes.title.text = text
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _xlsx_bytes() -> bytes:
    workbook = Workbook()
    workbook.active.title = "产品"
    workbook.active.append(["名称", "说明"])
    workbook.active.append(["数智名片", "企业展示"])
    buffer = io.BytesIO()
    workbook.save(buffer)
    return buffer.getvalue()


def test_csv_import_produces_drafts_and_rejects_formula_injection() -> None:
    payload = "title,raw_text,visibility\n产品说明,正文内容,internal\n".encode()
    assert validate_upload("docs.csv", "text/csv", payload) == "csv"
    assert parse_payload("csv", "docs.csv", payload)[0].visibility == "internal"

    with pytest.raises(KnowledgeImportError, match="IMPORT_DANGEROUS_VALUE"):
        parse_payload("csv", "docs.csv", b"title,raw_text\n=cmd,content\n")


def test_csv_import_accepts_tabular_columns_and_rejects_oversized_cells() -> None:
    draft = parse_payload("csv", "docs.csv", "title,raw_text,category\nA,content,产品\n".encode())[
        0
    ]
    assert "category: 产品" in draft.raw_text

    oversized = b"raw_text\n" + (b"a" * 100_001) + b"\n"
    with pytest.raises(KnowledgeImportError, match="IMPORT_CSV_CELL_TOO_LARGE"):
        parse_payload("csv", "docs.csv", oversized)


def test_docx_is_parsed_but_macro_and_archive_paths_are_rejected() -> None:
    payload = _docx_bytes("企业安全知识")
    assert (
        validate_upload(
            "guide.docx",
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            payload,
        )
        == "docx"
    )
    assert parse_payload("docx", "guide.docx", payload)[0].raw_text == "企业安全知识"

    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, "w") as archive:
        archive.writestr("word/document.xml", "<document/>")
        archive.writestr("../vbaProject.bin", b"macro")
    with pytest.raises(KnowledgeImportError, match="IMPORT_ARCHIVE_PATH"):
        validate_upload(
            "bad.docx",
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            buffer.getvalue(),
        )


def test_encrypted_pdf_and_mime_magic_mismatches_are_rejected() -> None:
    writer = PdfWriter()
    writer.add_blank_page(width=100, height=100)
    writer.encrypt("secret")
    buffer = io.BytesIO()
    writer.write(buffer)
    with pytest.raises(KnowledgeImportError, match="IMPORT_ENCRYPTED_PDF"):
        parse_payload("pdf", "secret.pdf", buffer.getvalue())

    with pytest.raises(KnowledgeImportError, match="IMPORT_MIME_MISMATCH"):
        validate_upload("file.pdf", "text/plain", b"%PDF-1.7")
    with pytest.raises(KnowledgeImportError, match="IMPORT_MAGIC_MISMATCH"):
        validate_upload("file.pdf", "application/pdf", b"not-a-pdf")


def test_pdf_uses_pymupdf_fallback_when_pypdf_cannot_recover_document(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    monkeypatch.setattr(
        knowledge_import_module,
        "PdfReader",
        lambda *_args, **_kwargs: (_ for _ in ()).throw(ValueError("broken xref")),
    )
    monkeypatch.setattr(
        knowledge_import_module,
        "_extract_pdf_text_with_fitz",
        lambda _payload: ("fallback text " * 10, 1),
    )

    draft = parse_payload("pdf", "fallback.pdf", b"%PDF-1.7 malformed")[0]

    assert draft.raw_text.startswith("fallback text")


def test_pdf_ocr_keeps_text_from_pages_after_a_blank_page(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    class _Pixmap:
        def __init__(self, page_number: int) -> None:
            self.page_number = page_number

        def tobytes(self, _format: str) -> bytes:
            return str(self.page_number).encode()

    class _Page:
        def __init__(self, page_number: int) -> None:
            self.page_number = page_number

        def get_pixmap(self, **_kwargs: object) -> _Pixmap:
            return _Pixmap(self.page_number)

    class _Document:
        page_count = 2

        def load_page(self, page_number: int) -> _Page:
            return _Page(page_number)

        def close(self) -> None:
            pass

    fake_fitz = SimpleNamespace(
        open=lambda **_kwargs: _Document(),
        Matrix=lambda *_args: object(),
    )
    monkeypatch.setitem(sys.modules, "fitz", fake_fitz)
    calls = iter([KnowledgeImportError("IMPORT_OCR_EMPTY"), "识别出的第二页文本"])

    def fake_ocr(_payload: bytes) -> str:
        result = next(calls)
        if isinstance(result, KnowledgeImportError):
            raise result
        return result

    monkeypatch.setattr(knowledge_import_module, "_ocr_image", fake_ocr)

    assert knowledge_import_module._ocr_pdf(b"%PDF-1.7", max_pages=2) == "识别出的第二页文本"


def test_office_and_html_formats_extract_text_without_network_access() -> None:
    pptx = _pptx_bytes("企业介绍")
    assert (
        validate_upload(
            "intro.pptx",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            pptx,
        )
        == "pptx"
    )
    assert "企业介绍" in parse_payload("pptx", "intro.pptx", pptx)[0].raw_text

    xlsx = _xlsx_bytes()
    assert (
        validate_upload(
            "catalog.xlsx",
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            xlsx,
        )
        == "xlsx"
    )
    assert "数智名片" in parse_payload("xlsx", "catalog.xlsx", xlsx)[0].raw_text

    html = (
        b"<h1>\xe4\xbc\x81\xe4\xb8\x9a</h1><script>ignore()</script><p>\xe4\xba\xa7\xe5\x93\x81</p>"
    )
    assert "企业" in parse_payload("html", "page.html", html)[0].raw_text
    assert "产品" in parse_payload("html", "page.html", html)[0].raw_text
    assert "ignore" not in parse_payload("html", "page.html", html)[0].raw_text


@pytest.mark.parametrize("name", ["../file.csv", "folder/file.csv", "folder\\file.csv", ""])
def test_unsafe_file_names_are_rejected(name: str) -> None:
    with pytest.raises(KnowledgeImportError, match="IMPORT_UNSAFE_FILENAME"):
        safe_file_name(name)
